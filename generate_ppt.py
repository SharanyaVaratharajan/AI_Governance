from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import os

# ---------------------------------------------------------
# Ensure docs folder exists
# ---------------------------------------------------------
output_dir = r"C:\Users\sharu\Documents\ai-governance-tower\docs"
os.makedirs(output_dir, exist_ok=True)

# ---------------------------------------------------------
# Helper: Create colorful diagram PNGs
# ---------------------------------------------------------
def create_diagram(filename, title, boxes):
    """
    boxes = list of tuples: (text, color)
    """
    img = Image.new("RGB", (1400, 800), "white")
    draw = ImageDraw.Draw(img)

    font = ImageFont.load_default()

    # Title
    draw.text((50, 30), title, fill="black", font=font)

    # Draw boxes
    x = 100
    y = 150
    w = 300
    h = 120
    spacing = 80

    for text, color in boxes:
        draw.rounded_rectangle((x, y, x + w, y + h), radius=25, fill=color)
        draw.text((x + 20, y + 45), text, fill="white", font=font)
        y += h + spacing

    img.save(filename)


# ---------------------------------------------------------
# Generate diagrams
# ---------------------------------------------------------
diagram1 = os.path.join(output_dir, "architecture.png")
create_diagram(
    diagram1,
    "System Architecture",
    [
        ("User Input", "#4A90E2"),
        ("FastAPI Gateway", "#50E3C2"),
        ("Governance Engine", "#9013FE"),
        ("SQLite Database", "#F5A623"),
        ("Dashboard UI", "#D0021B"),
    ],
)

diagram2 = os.path.join(output_dir, "governance_pipeline.png")
create_diagram(
    diagram2,
    "Governance Pipeline",
    [
        ("Input Payload", "#4A90E2"),
        ("Safety Flags", "#F8E71C"),
        ("Risk Scoring", "#7ED321"),
        ("Incident Creation", "#D0021B"),
        ("Audit Logging", "#9013FE"),
    ],
)

diagram3 = os.path.join(output_dir, "metadata_scan.png")
create_diagram(
    diagram3,
    "Column-Level Metadata Scan",
    [
        ("Column Name", "#4A90E2"),
        ("Rule Engine", "#50E3C2"),
        ("PII Detection", "#F5A623"),
        ("Tagging", "#9013FE"),
        ("Dashboard View", "#D0021B"),
    ],
)

# ---------------------------------------------------------
# Create PPT
# ---------------------------------------------------------
prs = Presentation()

def add_slide(title, bullets=None, image_path=None):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    if bullets:
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.level = 0

    if image_path:
        left = Inches(1)
        top = Inches(2.5)
        slide.shapes.add_picture(image_path, left, top, width=Inches(8))


# ---------------------------------------------------------
# Slides
# ---------------------------------------------------------

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "AI Governance Control Tower – Technical Deep Dive"
slide.placeholders[1].text = "Colorful Product Demo Style"

# Architecture Slide + Diagram
add_slide(
    "System Architecture",
    ["FastAPI", "Governance Engine", "SQLite", "Dashboard"],
    image_path=diagram1,
)

# Governance Pipeline + Diagram
add_slide(
    "Governance Pipeline",
    ["Input → Flags → Risk → Incident → Audit"],
    image_path=diagram2,
)

# Metadata Scan + Diagram
add_slide(
    "Column-Level Metadata Scan",
    ["PII detection", "Tagging", "Storage", "Dashboard"],
    image_path=diagram3,
)

# Other slides (text only)
add_slide("Backend (FastAPI)", ["Endpoints", "Gateway", "Mock Mode", "Risk Logic"])
add_slide("Dashboard", ["Systems", "Runs", "Incidents", "Heatmap", "Charts"])
add_slide("Database", ["Systems", "Runs", "Incidents", "Metadata"])
add_slide("Mock Mode", ["Offline", "Simulated Output", "Simulated Flags"])
add_slide("OpenAI Mode", ["Real Model Call", "Requires Credits"])
add_slide("Risk Scoring", ["Low", "Medium", "High"])
add_slide("Incident Creation", ["Triggers", "Severity", "Timeline"])
add_slide("Testing Workflow", ["Register", "Run", "Dashboard", "Incidents"])
add_slide("Troubleshooting", ["No runs", "No incidents", "Charts empty"])
add_slide("Future Enhancements", ["Lineage", "Compliance", "Risk Heatmaps"])

# ---------------------------------------------------------
# Save PPT
# ---------------------------------------------------------
output_path = os.path.join(output_dir, "governance_deep_dive_visual.pptx")
prs.save(output_path)

print("PPT saved to:", output_path)
